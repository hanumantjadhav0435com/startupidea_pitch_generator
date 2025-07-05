import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import logging

def create_pitch_ppt(pitch_data, startup_idea):
    """
    Create a PowerPoint presentation from pitch data
    """
    try:
        # Create presentation
        prs = Presentation()
        
        # Set slide size (16:9 aspect ratio)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        create_title_slide(prs, startup_idea, pitch_data['pitch'])
        
        # Create slides from pitch data
        for i, slide_data in enumerate(pitch_data['slides'], 1):
            create_content_slide(prs, slide_data['title'], slide_data['points'], i)
        
        # Add summary slide
        create_summary_slide(prs, pitch_data)
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_file.name)
        
        logging.info(f"Created PowerPoint presentation: {temp_file.name}")
        return temp_file.name
        
    except Exception as e:
        logging.error(f"Error creating PowerPoint: {str(e)}")
        raise Exception(f"Failed to create PowerPoint presentation: {str(e)}")

def create_title_slide(prs, startup_idea, pitch_line):
    """Create the title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = startup_idea
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Format title
    title_font = title.text_frame.paragraphs[0].runs[0].font
    title_font.size = Pt(44)
    title_font.bold = True
    title_font.color.rgb = RGBColor(31, 78, 121)
    
    # Subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = pitch_line
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Format subtitle
    subtitle_font = subtitle.text_frame.paragraphs[0].runs[0].font
    subtitle_font.size = Pt(24)
    subtitle_font.italic = True
    subtitle_font.color.rgb = RGBColor(68, 114, 196)

def create_content_slide(prs, title, points, slide_number):
    """Create a content slide with title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = f"{slide_number}. {title}"
    
    # Format title
    title_font = title_shape.text_frame.paragraphs[0].runs[0].font
    title_font.size = Pt(36)
    title_font.bold = True
    title_font.color.rgb = RGBColor(31, 78, 121)
    
    # Content
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    # Add bullet points
    for i, point in enumerate(points):
        p = text_frame.paragraphs[i] if i == 0 else text_frame.add_paragraph()
        p.text = point
        p.level = 0
        
        # Format bullet point
        font = p.runs[0].font
        font.size = Pt(20)
        font.color.rgb = RGBColor(68, 68, 68)
        
        # Add spacing between points
        if i < len(points) - 1:
            p.space_after = Pt(12)

def create_summary_slide(prs, pitch_data):
    """Create a summary slide with key information"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = "Key Takeaways"
    
    # Format title
    title_font = title_shape.text_frame.paragraphs[0].runs[0].font
    title_font.size = Pt(36)
    title_font.bold = True
    title_font.color.rgb = RGBColor(31, 78, 121)
    
    # Content
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    # Add key points
    key_points = [
        f"🎯 Problem: {pitch_data['problem'][:100]}...",
        f"💡 Solution: {pitch_data['solution'][:100]}...",
        f"📊 Market: {pitch_data['market'][:100]}...",
        f"💰 Revenue: {pitch_data['monetization'][:100]}..."
    ]
    
    for i, point in enumerate(key_points):
        p = text_frame.paragraphs[i] if i == 0 else text_frame.add_paragraph()
        p.text = point
        p.level = 0
        
        # Format point
        font = p.runs[0].font
        font.size = Pt(18)
        font.color.rgb = RGBColor(68, 68, 68)
        
        # Add spacing
        if i < len(key_points) - 1:
            p.space_after = Pt(12)

def cleanup_temp_file(file_path):
    """Clean up temporary PowerPoint file"""
    try:
        if os.path.exists(file_path):
            os.unlink(file_path)
            logging.info(f"Cleaned up temporary file: {file_path}")
    except Exception as e:
        logging.error(f"Error cleaning up temporary file: {str(e)}")
